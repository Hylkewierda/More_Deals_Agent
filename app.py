import os
import re
import json
import asyncio
from datetime import datetime
from typing import List, Optional, AsyncGenerator

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

from anthropic import AsyncAnthropic
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from openpyxl import load_workbook
from dotenv import load_dotenv
import io

load_dotenv()

app = FastAPI(title="More Deals Agent")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── Knowledge base ──────────────────────────────────────────────────────────

DOCS_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(DOCS_DIR, "data")
os.makedirs(DATA_DIR, exist_ok=True)

KNOWLEDGE_FILES = [
    ("MORE DEALS WEEK 1 - SALES FUNDAMENTALS.pptx", "Week 1 – Sales Fundamentals"),
    ("MORE DEALS WEEK 2 - FUNNEL BUILDING.pptx", "Week 2 – Funnel Building"),
    ("MORE DEALS WEEK 3 - SELLING.pptx", "Week 3 – Selling"),
    ("MORE DEALS WEEK 4 - NEGOTIATING.pptx", "Week 4 – Negotiating"),
    ("MORE DEALS WEEK 5 - CLOSING.pptx", "Week 5 – Closing"),
    ("sales_workshop.docx", "Sales Workshop – Verkopen & Onderhandelen voor Startups"),
    ("sales_workshop_2.docx", "Sales Workshop 2 – Eerste Salesgesprekken & Lead Generatie"),
    ("negotiation_workshop.docx", "Onderhandelen Workshop – Week 4"),
    ("fundraising_workshop.docx", "Fundraising Workshop – Fundraising voor Startups"),
    ("sales_session_structured.xlsx", "Sales Sessie – Gestructureerde Notities"),
]


def _read_docx(filepath: str) -> str:
    doc = Document(filepath)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_xlsx(filepath: str) -> str:
    wb = load_workbook(filepath, read_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        lines.append(f"[Sheet: {sheet}]")
        headers = [str(c) if c is not None else "" for c in rows[0]]
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                entry = " | ".join(f"{h}: {v}" for h, v in zip(headers, cells) if v.strip())
                if entry:
                    lines.append(entry)
        lines.append("")
    wb.close()
    return "\n".join(lines)


def _read_pptx(filepath: str) -> str:
    prs = Presentation(filepath)
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            lines.append(f"[Slide {slide_num}]")
            lines.extend(slide_texts)
            lines.append("")
    return "\n".join(lines)


SKILLS_DIR = os.path.join(DOCS_DIR, "skills")


def load_skills() -> str:
    """Load all skill files from data/skills/ directory."""
    if not os.path.exists(SKILLS_DIR):
        return ""
    sections = []
    for filename in sorted(os.listdir(SKILLS_DIR)):
        if filename.endswith(".md"):
            filepath = os.path.join(SKILLS_DIR, filename)
            try:
                with open(filepath, "r", encoding="utf-8") as f:
                    content = f.read()
                sections.append(content)
            except Exception as e:
                print(f"WARN: Skill laden mislukt {filename}: {e}")
    return "\n\n---\n\n".join(sections)


CACHE_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data", "knowledge_cache.json")


def _cache_is_valid(cache: dict) -> bool:
    for filename, _ in KNOWLEDGE_FILES:
        filepath = os.path.join(DOCS_DIR, filename)
        if not os.path.exists(filepath):
            continue
        mtime = os.path.getmtime(filepath)
        if cache.get("mtimes", {}).get(filename) != mtime:
            return False
    return True


def load_knowledge_base() -> str:
    os.makedirs(os.path.dirname(CACHE_FILE), exist_ok=True)

    # Load from cache if valid
    if os.path.exists(CACHE_FILE):
        try:
            with open(CACHE_FILE, "r", encoding="utf-8") as f:
                cache = json.load(f)
            if _cache_is_valid(cache):
                print("Kennisbasis geladen vanuit cache.")
                return cache["content"]
        except Exception:
            pass

    # Parse source files and build cache
    sections = []
    mtimes = {}
    for filename, title in KNOWLEDGE_FILES:
        filepath = os.path.join(DOCS_DIR, filename)
        if not os.path.exists(filepath):
            print(f"WARN: Niet gevonden: {filename}")
            continue
        try:
            if filename.lower().endswith(".pptx"):
                text = _read_pptx(filepath)
            elif filename.lower().endswith(".xlsx"):
                text = _read_xlsx(filepath)
            else:
                text = _read_docx(filepath)
            sections.append(f"### {title}\n\n{text}")
            mtimes[filename] = os.path.getmtime(filepath)
        except Exception as e:
            print(f"ERROR: Fout bij laden {filename}: {e}")

    content = "\n\n---\n\n".join(sections)

    # Save cache
    try:
        with open(CACHE_FILE, "w", encoding="utf-8") as f:
            json.dump({"content": content, "mtimes": mtimes}, f, ensure_ascii=False, indent=2)
        print(f"Kennisbasis gecached naar {CACHE_FILE}")
    except Exception as e:
        print(f"WARN: Cache opslaan mislukt: {e}")

    return content


print("Kennisbasis laden…")
KNOWLEDGE_BASE = load_knowledge_base()
print(f"Kennisbasis geladen: {len(KNOWLEDGE_BASE):,} tekens")

print("Skills laden…")
SKILLS_CONTENT = load_skills()
print(f"Skills geladen: {len(SKILLS_CONTENT):,} tekens")

SYSTEM_INSTRUCTIONS = """Je bent de **More Deals AI Agent** — een expert sales coach en assistent van het More Deals programma, opgericht door Joey Moreau.

More Deals helpt ondernemers en startups om hun salescijfers te verhogen via een gestructureerd programma dat bestaat uit:
- Week 1: Sales Fundamentals
- Week 2: Funnel Building
- Week 3: Selling (LEADS framework)
- Week 4: Negotiating
- Week 5: Closing

**Jouw rol:**
- Beantwoord vragen over sales, onderhandelen, lead generatie en fundraising
- Geef praktisch, concreet advies op basis van het More Deals programma
- Leg frameworks uit (LEADS, MAGIC, Feedback Loop, Certainty Script, etc.)
- Help bij het kwalificeren van leads en deals
- Schrijf outreach berichten, follow-ups en reminders met behulp van je outreach skills
- Spreek Nederlands tenzij de gebruiker Engels spreekt

**Communicatiestijl:**
- Direct en to-the-point, zoals Joey Moreau dat doet
- Gebruik concrete voorbeelden en scripts
- Geef altijd actionable stappen
- Wees enthousiast maar professioneel
- Gebruik de frameworks en taal uit de trainingen

**Formatting regels:**
- Gebruik NOOIT emoji's in je antwoorden. Geen enkele emoji, nergens.
- Structureer antwoorden duidelijk met markdown: koppen (##, ###), bullet points, genummerde lijsten en **vetgedrukte** termen
- Gebruik witregels tussen secties voor leesbaarheid
- Houd alinea's kort en scanbaar (max 3-4 zinnen per alinea)
- Zet frameworks, methodes en belangrijke termen altijd **vetgedrukt**

Hieronder staat je volledige kennisbasis. Gebruik deze als bron voor al je antwoorden."""

SKILLS_SECTION = f"\n\n## OUTREACH SKILLS\n\nHieronder staan je outreach skills. Gebruik deze wanneer een gebruiker vraagt om outreach berichten, cold messages, follow-ups, tweede berichten na een reactie, of reminders te schrijven. Volg de structuur, regels en het proces uit de betreffende skill exact op.\n\n{SKILLS_CONTENT}" if SKILLS_CONTENT else ""

SYSTEM_PROMPT_WITH_KB = f"{SYSTEM_INSTRUCTIONS}{SKILLS_SECTION}\n\n## KENNISBASIS – MORE DEALS PROGRAMMA\n\n{KNOWLEDGE_BASE}"

# ── Anthropic client ─────────────────────────────────────────────────────────

MODEL = "claude-haiku-4-5-20251001"

def get_client() -> AsyncAnthropic:
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(
            status_code=500,
            detail="ANTHROPIC_API_KEY niet ingesteld. Voeg hem toe aan je .env bestand."
        )
    return AsyncAnthropic(api_key=api_key)


# ── Pydantic models ─────────────────────────────────────────────────────────

class Message(BaseModel):
    role: str  # "user" | "assistant"
    content: str


class ChatRequest(BaseModel):
    messages: List[Message]


class LeadForm(BaseModel):
    naam: str
    bedrijf: str
    email: str
    telefoon: Optional[str] = None
    sector: Optional[str] = None
    bedrijfsgrootte: Optional[str] = None
    pijnpunten: str
    budget: Optional[str] = None
    tijdlijn: Optional[str] = None
    notities: Optional[str] = None


class DealForm(BaseModel):
    prospect_naam: str
    bedrijf: str
    email: Optional[str] = None
    deal_waarde: Optional[str] = None
    fase: str  # prospecting, gekwalificeerd, voorstel, onderhandeling, closing, gewonnen, verloren
    volgende_stap: str
    follow_up_datum: Optional[str] = None
    notities: Optional[str] = None


# ── Query routing ──────────────────────────────────────────────────────────

SALES_KEYWORDS = re.compile(
    r"\b("
    # Kern sales
    r"sales|lead|leads|funnel|pitch|deal|deals|onderhandel|closing|close"
    r"|prospect|follow.?up|offerte|klant|klanten|verkoop|verkopen|more deals"
    # Frameworks & methodes
    r"|leads framework|magic|certainty script|feedback loop"
    r"|SPIN|MEDDIC|challenger|consultative|value.?selling"
    # Bezwaren & prijs
    r"|bezwaar|bezwaren|twijfel|weerstand|afwijzing|overtuig"
    r"|prijs|korting|marge|pricing|investering|kosten|budget|pakket|abonnement"
    # Pipeline & metrics
    r"|pipeline|kwalificer|acquisitie|cold.?call|cold.?mail|outreach"
    r"|voorstel|contract|omzet|target|quota|conversie|forecast|hitrate|win.?rate"
    r"|revenue|MRR|ARR|churn|LTV|CAC|ROI|groei|winst|marktaandeel"
    # Prospect & klantrelatie
    r"|doelgroep|segment|ICP|account|relatie|stakeholder|beslisser|beslissing|DMU"
    r"|referral|referentie|aanbeveling|ambassadeur|testimonial|case.?study"
    r"|retentie|loyaliteit|upsell|cross.?sell|nurture|reactiveer"
    # Communicatie & tools
    r"|demo|presentatie|afspraak|meeting|kennismak|introductie|connectie"
    r"|LinkedIn|script|belscript|template|openingszin|campagne|webinar"
    r"|CRM|salesforce|hubspot|database|whitepaper|beurs|event"
    r"|bellen|opvolg|terugbel|outreach|eerste bericht|second message|tweede bericht"
    r"|follow.?up bericht|reminder|herinnering|cold.?message|linkedin bericht"
    # Waarde & propositie
    r"|waarde|waardepropositie|propositie|USP|pijnpunt|oplossing|resultaat"
    r"|concurrentie|markt|vertrouwen|akkoord|handtekening|commitment|urgentie"
    # Mindset & training
    r"|mindset|motivatie|discipline|role.?play|coaching"
    # Bedrijfstypen
    r"|B2B|B2C|enterprise|MKB|ZZP|scale.?up|ondernemer|opdracht"
    # Fundraising
    r"|fundrais|investor|investeerder|startup|seed|raise"
    # Cursus
    r"|week [1-5]|workshop"
    r")\b",
    re.IGNORECASE,
)


def needs_knowledge_base(query: str) -> bool:
    return bool(SALES_KEYWORDS.search(query)) or "DOCUMENT_EXCERPT" in query


# ── Chat endpoint (streaming) ───────────────────────────────────────────────

@app.post("/api/chat")
async def chat(request: ChatRequest):
    client = get_client()

    messages = [{"role": m.role, "content": m.content} for m in request.messages]

    last_user_msg = next((m.content for m in reversed(request.messages) if m.role == "user"), "")
    use_kb = needs_knowledge_base(last_user_msg)
    system_text = SYSTEM_PROMPT_WITH_KB if use_kb else f"{SYSTEM_INSTRUCTIONS}{SKILLS_SECTION}"

    async def generate() -> AsyncGenerator[str, None]:
        try:
            if use_kb:
                if "DOCUMENT_EXCERPT" in last_user_msg:
                    steps = [
                        "Document verwerken...",
                        "Inhoud analyseren...",
                        "Antwoord formuleren...",
                    ]
                else:
                    steps = [
                        "Vraag analyseren...",
                        "Relevante context ophalen...",
                        "Antwoord formuleren...",
                    ]
                for step in steps:
                    yield f"data: {json.dumps({'status': step})}\n\n"
                    await asyncio.sleep(0.6)

            async with client.messages.stream(
                model=MODEL,
                max_tokens=2048,
                system=system_text,
                messages=messages,
            ) as stream:
                async for text in stream.text_stream:
                    if text:
                        yield f"data: {json.dumps({'text': text})}\n\n"
            yield "data: [DONE]\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'text': f'**API-fout:** {str(e)}'})}\n\n"
            yield "data: [DONE]\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


# ── File Extraction Endpoint ────────────────────────────────────────────────

@app.post("/api/extract")
async def extract_text(file: UploadFile = File(...)):
    name = file.filename.lower()
    content = await file.read()
    text = ""
    try:
        if name.endswith(".docx"):
            doc = Document(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif name.endswith(".pdf"):
            reader = PdfReader(io.BytesIO(content))
            pages = [page.extract_text() or "" for page in reader.pages]
            text = "\n\n".join(p for p in pages if p.strip())
        elif name.endswith((".txt", ".md", ".csv", ".json")):
            text = content.decode("utf-8")
        else:
            return {"error": "Unsupported file type. Please use .pdf, .txt, .md, .csv, .json, or .docx."}
    except Exception as e:
        return {"error": f"Failed to parse file: {str(e)}"}
    
    return {"text": text, "filename": file.filename}

# ── Lead endpoints ──────────────────────────────────────────────────────────

def load_json(filename: str) -> list:
    path = os.path.join(DATA_DIR, filename)
    if not os.path.exists(path):
        return []
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def save_json(filename: str, data: list) -> None:
    path = os.path.join(DATA_DIR, filename)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


@app.post("/api/leads")
async def create_lead(lead: LeadForm):
    leads = load_json("leads.json")
    entry = lead.model_dump()
    entry["id"] = len(leads) + 1
    entry["aangemaakt_op"] = datetime.now().isoformat()
    leads.append(entry)
    save_json("leads.json", leads)
    return {"success": True, "id": entry["id"], "bericht": f"Lead '{lead.naam}' succesvol opgeslagen!"}


@app.get("/api/leads")
async def get_leads():
    return load_json("leads.json")


@app.post("/api/deals")
async def create_deal(deal: DealForm):
    deals = load_json("deals.json")
    entry = deal.model_dump()
    entry["id"] = len(deals) + 1
    entry["aangemaakt_op"] = datetime.now().isoformat()
    deals.append(entry)
    save_json("deals.json", deals)
    return {"success": True, "id": entry["id"], "bericht": f"Deal '{deal.prospect_naam}' succesvol opgeslagen!"}


@app.get("/api/deals")
async def get_deals():
    return load_json("deals.json")


# ── Conversation endpoints ──────────────────────────────────────────────────

class ConversationSave(BaseModel):
    title: Optional[str] = None
    messages: List[Message]


@app.post("/api/conversations")
async def save_conversation(conv: ConversationSave):
    conversations = load_json("conversations.json")
    first_user = next((m.content for m in conv.messages if m.role == "user"), "Nieuw gesprek")
    title = conv.title or (first_user[:60] + ("..." if len(first_user) > 60 else ""))
    entry = {
        "id": datetime.now().strftime("%Y%m%d%H%M%S%f"),
        "title": title,
        "saved_at": datetime.now().isoformat(),
        "messages": [{"role": m.role, "content": m.content} for m in conv.messages],
    }
    conversations.append(entry)
    save_json("conversations.json", conversations)
    return {"success": True, "id": entry["id"], "title": entry["title"]}


@app.get("/api/conversations")
async def list_conversations():
    convs = load_json("conversations.json")
    return [
        {"id": c["id"], "title": c["title"], "saved_at": c["saved_at"], "message_count": len(c["messages"])}
        for c in reversed(convs)
    ]


@app.get("/api/conversations/{conv_id}")
async def get_conversation(conv_id: str):
    conversations = load_json("conversations.json")
    conv = next((c for c in conversations if c["id"] == conv_id), None)
    if not conv:
        raise HTTPException(status_code=404, detail="Gesprek niet gevonden")
    return conv


@app.delete("/api/conversations/{conv_id}")
async def delete_conversation(conv_id: str):
    conversations = load_json("conversations.json")
    conversations = [c for c in conversations if c["id"] != conv_id]
    save_json("conversations.json", conversations)
    return {"success": True}


# ── Static files & root ──────────────────────────────────────────────────────

app.mount("/static", StaticFiles(directory=os.path.join(DOCS_DIR, "static")), name="static")


@app.get("/")
async def root():
    response = FileResponse(os.path.join(DOCS_DIR, "static", "index.html"))
    response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
    response.headers["Pragma"] = "no-cache"
    response.headers["Expires"] = "0"
    return response


@app.get("/health")
async def health():
    return {"status": "ok", "knowledge_chars": len(KNOWLEDGE_BASE)}
